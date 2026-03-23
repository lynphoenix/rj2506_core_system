import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. 初始化 PPT 对象
prs = Presentation()

# 自定义主题颜色
COLOR_PRIMARY = RGBColor(0, 51, 102)     # 深蓝
COLOR_SECONDARY = RGBColor(0, 102, 204)   # 科技蓝
COLOR_ACCENT = RGBColor(255, 102, 0)      # 橙色点缀
COLOR_TEXT_DARK = RGBColor(51, 51, 51)    # 深灰文字
COLOR_TEXT_LIGHT = RGBColor(255, 255, 255) # 白色文字
COLOR_BG_LIGHT = RGBColor(245, 247, 250)  # 浅灰蓝背景

# 设置全局字体
def set_font(p, name='Microsoft YaHei'):
    p.font.name = name

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_LIGHT

def add_header_footer(slide, title_text):
    # 添加顶部装饰条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_PRIMARY
    top_bar.line.fill.background()

    # 添加页脚
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(9), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "RJ2506 双臂轮式复合机器人 | 核心技术攻坚与系统架构"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(120, 120, 120)
    set_font(p)

# --- Slide 1: 封面 ---
slide = prs.slides.add_slide(prs.slide_layouts[6]) # 空白版式
set_slide_background(slide)

# 装饰背景色块
bg_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(10), Inches(3)
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = COLOR_PRIMARY
bg_shape.line.fill.background()

# 主标题
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "RJ2506 全局架构与任务拆解"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = COLOR_TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

# 副标题
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "双臂轮式复合机器人系统方案 V3.0\n(新增顶层大脑节点与完整 8 段式管线)"
p2.font.size = Pt(26)
p2.font.color.rgb = RGBColor(200, 220, 255)
p2.alignment = PP_ALIGN.CENTER

# 汇报人/日期
txBox3 = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Project: Aloha / Date: 2026.03"
p3.font.size = Pt(16)
p3.font.bold = True
p3.font.color.rgb = COLOR_TEXT_DARK
p3.alignment = PP_ALIGN.CENTER
set_font(p3)


# --- Slide 2: 整体任务描述与核心需求 ---
slide = prs.slides.add_slide(prs.slide_layouts[5]) # 仅标题
set_slide_background(slide)
add_header_footer(slide, "项目背景")

title = slide.shapes.title
title.text = "整体任务描述 (The Mission)"
p_title = title.text_frame.paragraphs[0]
p_title.font.color.rgb = COLOR_PRIMARY
p_title.font.bold = True
set_font(p_title)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "RJ2506 是一款双臂轮式复合机器人，旨在解决工厂半结构化环境下的高精度柔性装配、堆垛与搬运问题。"
p.font.size = Pt(22)
p.font.color.rgb = COLOR_TEXT_DARK

# 添加两个卡片展示核心场景
# 场景 1 卡片
card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.8), Inches(4.2), Inches(3.8))
card1.fill.solid()
card1.fill.fore_color.rgb = RGBColor(255, 255, 255)
card1.line.color.rgb = COLOR_SECONDARY

tf1 = card1.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.text = "场景 1：大框码料"
p1.font.bold = True
p1.font.size = Pt(22)
p1.font.color.rgb = COLOR_SECONDARY

p1_desc = tf1.add_paragraph()
p1_desc.text = """• 目标: 双手从下料平台同时抓取大料片，搬运并精准码放到大框中
• 难点: 大框边缘限制，放置必须极度平稳精准，杜绝磕碰"""
p1_desc.font.size = Pt(16)
p1_desc.font.color.rgb = COLOR_TEXT_DARK

# 场景 2 卡片
card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(2.8), Inches(4.2), Inches(3.8))
card2.fill.solid()
card2.fill.fore_color.rgb = RGBColor(255, 255, 255)
card2.line.color.rgb = COLOR_ACCENT

tf2 = card2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "场景 2：小框换框与码垛"
p2.font.bold = True
p2.font.size = Pt(22)
p2.font.color.rgb = COLOR_ACCENT

p2_desc = tf2.add_paragraph()
p2_desc.text = """• 目标: 推动空框替换满框，并将满框双手搬运至中转站进行高精度码垛
• 难点: 视线完全遮挡，码垛极易重心偏移倒塌，需抗盲插和倾斜修正"""
p2_desc.font.size = Pt(16)
p2_desc.font.color.rgb = COLOR_TEXT_DARK


# --- Slide 3: 硬件载体与计算平台 ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_slide_background(slide)
add_header_footer(slide, "硬件基础")

title = slide.shapes.title
title.text = "硬件载体与底层通信"
p_title = title.text_frame.paragraphs[0]
p_title.font.color.rgb = COLOR_PRIMARY
p_title.font.bold = True
set_font(p_title)

left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(4)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True

def add_bullet(tf, text, level=0, bold=False, color=COLOR_TEXT_DARK):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    # 增加基础字号，提升可读性
    p.font.size = Pt(28 - level*4)
    p.font.bold = bold
    p.font.color.rgb = color
    set_font(p)
    # 设置行距和段前距
    p.line_spacing = 1.2
    p.space_before = Pt(8)

# 清除默认段落
tf.clear()

add_bullet(tf, "核心硬件平台", bold=True, color=COLOR_PRIMARY)
add_bullet(tf, "RJ2506 差速底盘 + 躯干控制器", level=1)
add_bullet(tf, "双机械臂 (单臂负载 15kg，协同 30kg)", level=1)
add_bullet(tf, "左右平行夹爪", level=1)
add_bullet(tf, "头部全局相机(OAK-D) + 手腕局部相机", level=1)

add_bullet(tf, "计算平台与系统", bold=True, color=COLOR_PRIMARY)
add_bullet(tf, "NVIDIA Jetson Orin", level=1)
add_bullet(tf, "Phase 0/1 阶段暂不刷入 PREEMPT_RT 实时内核 (避开现有任务冲突)", level=1)

add_bullet(tf, "底层通信总线", bold=True, color=COLOR_ACCENT)
add_bullet(tf, "CANopen (1000Hz) - 控制 14 轴电机", level=1)
add_bullet(tf, "RS-485 (50Hz) - 控制双侧夹爪", level=1)


# --- Slide 4: 全局软件架构: 八段式控制管线 ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_slide_background(slide)
add_header_footer(slide, "核心架构")

title = slide.shapes.title
title.text = "全局软件架构：八段式控制管线"
p_title = title.text_frame.paragraphs[0]
p_title.font.color.rgb = COLOR_PRIMARY
p_title.font.bold = True
set_font(p_title)

# 绘制 8 个阶段的流程图卡片
stages = [
    "[0] 顶层大脑: 负责全程调度、流转与异常处理",
    "[1] 任务触发与备料: 接收信号，执行初期的备料等准备",
    "[2] 宏观导航: 负责设备移动",
    "[3] 物理顺应补偿: 利用算法控制消除底盘误差等对位不准",
    "[4] 大尺度抓取: 系统感知目标位姿并完成抱持",
    "[5] 高空特征记忆: 在目标上方悬停，预先提取环境特征",
    "[6] 微插闭环: 在视线遮挡阶段接管，进行姿态等参数修正",
    "[7] 触底检测与释放: 检测放置到位并释放",
    "[8] 复位与回程: 收回机械臂并返回"
]

# 脑节点跨越所有
brain_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(0.6))
brain_box.fill.solid()
brain_box.fill.fore_color.rgb = COLOR_PRIMARY
brain_box.text_frame.text = stages[0]
brain_box.text_frame.paragraphs[0].font.bold = True
brain_box.text_frame.paragraphs[0].font.size = Pt(16)
brain_box.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT_LIGHT

# 下方阶段分为两列排列
for i in range(1, 9):
    col = 0 if i <= 4 else 1
    row = (i - 1) % 4

    left = Inches(0.5 + col * 4.6)
    top = Inches(2.3 + row * 1.1)
    width = Inches(4.4)
    height = Inches(0.9)

    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()

    # 区分不同类型的阶段颜色
    if i in [2]: # 导航
        box.fill.fore_color.rgb = RGBColor(220, 235, 255)
    elif i in [4]: # AI
        box.fill.fore_color.rgb = RGBColor(255, 235, 210)
    elif i in [3, 7]: # 力控与交互
        box.fill.fore_color.rgb = RGBColor(255, 220, 220)
    elif i in [5, 6]: # 视觉闭环 (核心攻坚)
        box.fill.fore_color.rgb = RGBColor(220, 255, 220)
        box.line.color.rgb = RGBColor(0, 150, 0)
        box.line.width = Pt(2)
    else:
        box.fill.fore_color.rgb = RGBColor(240, 240, 240)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = stages[i]
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT_DARK

    # 如果不是最后一行，画个向下箭头
    if row < 3:
        slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left + Inches(2.1), top + Inches(0.9), Inches(0.2), Inches(0.2)).fill.solid()


# --- Slide 5: 架构演进与破局思路 ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_slide_background(slide)
add_header_footer(slide, "架构演进")

title = slide.shapes.title
title.text = "为什么采用这套管线设计？"
p_title = title.text_frame.paragraphs[0]
p_title.font.color.rgb = COLOR_PRIMARY
p_title.font.bold = True
set_font(p_title)

# 痛点框
box_pain = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(4.2), Inches(4.5))
box_pain.fill.solid()
box_pain.fill.fore_color.rgb = RGBColor(255, 230, 230) # 浅红
box_pain.line.color.rgb = RGBColor(200, 0, 0)
tf_pain = box_pain.text_frame
tf_pain.word_wrap = True
p_pain = tf_pain.paragraphs[0]
p_pain.text = "❌ 现有方案难点"
p_pain.font.bold = True
p_pain.font.size = Pt(26)
p_pain.font.color.rgb = RGBColor(200, 0, 0)

add_bullet(tf_pain, "", level=0)
add_bullet(tf_pain, "端到端模型控制边界", level=1, bold=True)
add_bullet(tf_pain, "难以确保导航和物理接触的绝对安全", level=2)
add_bullet(tf_pain, "底盘停靠天然误差", level=1, bold=True)
add_bullet(tf_pain, "对精密装配/码垛造成挑战", level=2)
add_bullet(tf_pain, "末端盲区", level=1, bold=True)
add_bullet(tf_pain, "大型料框遮挡相机视线", level=2)

# 箭头
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.8), Inches(3.8), Inches(0.4), Inches(0.5)).fill.solid()

# 破局框
box_sol = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(1.8), Inches(4.2), Inches(4.5))
box_sol.fill.solid()
box_sol.fill.fore_color.rgb = RGBColor(230, 245, 230) # 浅绿
box_sol.line.color.rgb = RGBColor(0, 150, 0)
tf_sol = box_sol.text_frame
tf_sol.word_wrap = True
p_sol = tf_sol.paragraphs[0]
p_sol.text = "✅ 我们的架构方案"
p_sol.font.bold = True
p_sol.font.size = Pt(26)
p_sol.font.color.rgb = RGBColor(0, 150, 0)

add_bullet(tf_sol, "", level=0)
add_bullet(tf_sol, "引入顶层大脑节点", level=1, bold=True)
add_bullet(tf_sol, "统一调度流转与异常处理", level=2)
add_bullet(tf_sol, "物理顺应补偿", level=1, bold=True)
add_bullet(tf_sol, "利用力控算法抹平底盘误差", level=2)
add_bullet(tf_sol, "高空特征记忆 + 盲插修正", level=1, bold=True)
add_bullet(tf_sol, "依靠预存特征进行姿态闭环", level=2)


# --- Slide 6: 场景 1 映射分析 ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_slide_background(slide)
add_header_footer(slide, "大框码料")

title = slide.shapes.title
title.text = "场景一：大框码料流程分析"
p_title = title.text_frame.paragraphs[0]
p_title.font.color.rgb = COLOR_PRIMARY
p_title.font.bold = True
set_font(p_title)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
tf.clear()

add_bullet(tf, "动作 1-3: 寻位与对齐", bold=True, color=COLOR_SECONDARY)
add_bullet(tf, "[1] 触发备料: 从待机休眠模式唤醒，确认双臂处于安全行车姿态", level=1)
add_bullet(tf, "[2] 宏观导航: 底盘导航并宽容停靠在冲压下料平台前", level=1)
add_bullet(tf, "[3] 物理顺应: 双手伸出抵住下料平台边缘，利用阻抗控制将自身位姿卡死对齐", level=1)

add_bullet(tf, "", level=0)
add_bullet(tf, "动作 4-6: 抓取与放置", bold=True, color=COLOR_SECONDARY)
add_bullet(tf, "[4] 大尺度抓取: 视觉输出协作轨迹，双手准确抓取 5-6 片大料片", level=1)
add_bullet(tf, "[5] 高空记忆: 转身移动到大框旁边，在大框上方15cm处悬停，记录边角特征", level=1)
add_bullet(tf, "[6] 微插闭环: 带着宽大的料片往下放，视线全盲，切入虚拟算法实时修正姿态", level=1)

add_bullet(tf, "", level=0)
add_bullet(tf, "动作 7-8: 完成与复位", bold=True, color=COLOR_SECONDARY)
add_bullet(tf, "[7] 触底释放: 监听Z轴电流突变感知放到底部，双侧夹爪同步松手", level=1)
add_bullet(tf, "[8] 复位回程: 双臂退回安全姿态，完成一次循环，系统复位", level=1)

# --- Slide 7: 场景 2 映射分析 ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_slide_background(slide)
add_header_footer(slide, "小框换框")

title = slide.shapes.title
title.text = "场景二：小框换框与码垛流程分析"
p_title = title.text_frame.paragraphs[0]
p_title.font.color.rgb = COLOR_PRIMARY
p_title.font.bold = True
set_font(p_title)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
tf.clear()

add_bullet(tf, "动作 1-3: 物理替换", bold=True, color=COLOR_SECONDARY)
add_bullet(tf, "[1] 触发备料: 唤醒后，指引底盘先前往空框区，抓取一个空框作为替换弹药", level=1)
add_bullet(tf, "[2] 宏观导航: 带着空框移动到冲压机台下料口停靠", level=1)
add_bullet(tf, "[3] 物理顺应: 用手里的空框顺着机台导轨往里推，顺势把满框顶出来", level=1)

add_bullet(tf, "", level=0)
add_bullet(tf, "动作 4-6: 搬运与高精度码垛", bold=True, color=COLOR_SECONDARY)
add_bullet(tf, "[4] 大尺度抓取: 模型感知顶出的满框位姿，双手从两侧牢牢抱起满框", level=1)
add_bullet(tf, "[5] 高空记忆: 搬运满框前往中转站，在当前堆叠层上方悬停提取下方榫卯特征", level=1)
add_bullet(tf, "[6] 微插闭环: 往下放置时满框挡住视线，依靠矩阵修正倾斜，防止 5层码垛倒塌", level=1)

add_bullet(tf, "", level=0)
add_bullet(tf, "动作 7-8: 完成与复位", bold=True, color=COLOR_SECONDARY)
add_bullet(tf, "[7] 触底释放: 接触力矩达到峰值判定码放稳固，松开双臂", level=1)
add_bullet(tf, "[8] 复位回程: 退回行车姿态，准备下一次换框任务", level=1)


# --- Slide 8: Phase 0 ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_slide_background(slide)
add_header_footer(slide, "开发计划")

title = slide.shapes.title
title.text = "系统开发与任务拆解: Phase 0"
p_title = title.text_frame.paragraphs[0]
p_title.font.color.rgb = COLOR_PRIMARY
p_title.font.bold = True
set_font(p_title)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
tf.clear()

add_bullet(tf, "Phase 0: 基础设施与物理真值采集", bold=True, color=COLOR_SECONDARY, level=0)
add_bullet(tf, "时间: 2周", level=1, color=COLOR_ACCENT, bold=True)
add_bullet(tf, "任务: ROS 2环境安装、相机外参标定、探针制作与Ground Truth采集、特征图金标准采集", level=1)
add_bullet(tf, "指标: 软时间插值抖动 < 10ms, 重投影误差 RMSE < 1.5像素", level=1)

# --- Slide 9: Phase 1 ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_slide_background(slide)
add_header_footer(slide, "开发计划")

title = slide.shapes.title
title.text = "系统开发与任务拆解: Phase 1"
p_title = title.text_frame.paragraphs[0]
p_title.font.color.rgb = COLOR_PRIMARY
p_title.font.bold = True
set_font(p_title)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
tf.clear()

add_bullet(tf, "Phase 1: 核心硬骨头算法模块攻坚 (研发重点)", bold=True, color=COLOR_SECONDARY, level=0)
add_bullet(tf, "时间: 6周", level=1, color=COLOR_ACCENT, bold=True)
add_bullet(tf, "任务:", level=1)
add_bullet(tf, "• 双臂全身协调控制(WBC)", level=2)
add_bullet(tf, "• 顶层状态机", level=2)
add_bullet(tf, "• 双目特征提取", level=2)
add_bullet(tf, "• 虚拟IBVS视觉伺服", level=2)
add_bullet(tf, "• C++守护进程、阻抗控制", level=2)

# --- Slide 10: Phase 2 ---
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_slide_background(slide)
add_header_footer(slide, "开发计划")

title = slide.shapes.title
title.text = "系统开发与任务拆解: Phase 2"
p_title = title.text_frame.paragraphs[0]
p_title.font.color.rgb = COLOR_PRIMARY
p_title.font.bold = True
set_font(p_title)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
tf = txBox.text_frame
tf.word_wrap = True
tf.clear()

add_bullet(tf, "Phase 2: AI 训练与全链路联调", bold=True, color=COLOR_SECONDARY, level=0)
add_bullet(tf, "时间: 4周", level=1, color=COLOR_ACCENT, bold=True)
add_bullet(tf, "任务: Sim2Real数据采集与训练、模型部署、实车打通与联调", level=1)


# --- Slide 11: 结尾 ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide)

# 渐变装饰
bg_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = COLOR_PRIMARY

txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "THANKS"
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = COLOR_TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "只有底层基石 100% 夯实，AI 端到端才有意义。"
p2.font.size = Pt(26)
p2.font.color.rgb = RGBColor(200, 220, 255)
p2.alignment = PP_ALIGN.CENTER

prs.save('docs/RJ2506_Architecture_V3.pptx')
print("PPTX generated successfully: docs/RJ2506_Architecture_V3.pptx")
